"""Update FINANO BOCHK Challenge 2026 PPT with competition-aligned content."""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


SRC = Path(
    r"c:\Users\30362\xwechat_files\wxid_iioe8mapqfa322_92bd\msg\file\2026-06\FINANO_Draft_0609.pptx"
)
OUT = SRC.with_name("FINANO_Draft_0609_BOCHK2026.pptx")
COPY = Path(r"D:\FINANO\FINANO_Draft_0609_BOCHK2026.pptx")


def set_text(shape, text: str) -> None:
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text


def find_textbox_in_group(group, index: int = 1):
    boxes = [s for s in group.shapes if hasattr(s, "text_frame")]
    if boxes:
        return boxes[min(index, len(boxes) - 1)]
    return None


def replace_in_slide(slide, old: str, new: str) -> int:
    count = 0
    for shape in slide.shapes:
        if hasattr(shape, "text") and old in shape.text:
            shape.text = shape.text.replace(old, new)
            count += 1
        if shape.shape_type == 6:
            for sub in shape.shapes:
                if hasattr(sub, "text") and old in sub.text:
                    sub.text = sub.text.replace(old, new)
                    count += 1
    return count


def main() -> None:
    prs = Presentation(str(SRC))

    # --- Slide 1: cover fields ---
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.shape_type == 6:
            inner = find_textbox_in_group(shape)
            if inner is None:
                continue
            if inner.text.strip() == "FINANO":
                set_text(inner, "FINANO")
            elif "Team Leader" in inner.text:
                set_text(
                    inner,
                    "Team Leader: Pu Tianyi (HKUST)\n"
                    "Wang Xinyi | Wang Zhenghao | Zhu Yijin\n"
                    "組長：浦天祎｜隊員：王心怡、王正昊、朱艺瑾",
                )
            elif "Student Stream" in inner.text or "學生組" in inner.text:
                set_text(inner, "學生組 Student Stream")

    # Proposal summary box (empty rounded rectangle)
    summary = (
        "FINANO is a Generative AI multi-agent engine for bank-side fund suitability. "
        "It adds a behavioral layer (FBTI) and MAFB parallel analysis with compliance "
        "logging on top of existing KYC—helping BOCHK reduce mis-selling risk and "
        "document every recommendation. Tech: GenAI + Big Data | Scenario: Risk Mgmt.\n\n"
        "FINANO 是以生成式 AI 與多智能體打造的銀行端基金適配引擎，在既有 KYC 之上"
        "加入行為畫像（FBTI）與合規留痕的 MAFB 分析，協助中銀香港提升適當性管理、"
        "降低違規銷售風險。科技主題：生成式 AI、大數據；場景主題：風控、普惠金融。"
    )
    for shape in s1.shapes:
        if hasattr(shape, "text_frame") and shape.text.strip() == "":
            set_text(shape, summary)
            break

    # --- Slide 2: competition themes (official template selections) ---
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if hasattr(shape, "text") and (
            "Quantum" in shape.text
            or "Tech Theme" in shape.text
            or "科技主題" in shape.text
        ):
            set_text(
                shape,
                "【已選科技主題 Selected Technology】\n"
                "✓ 生成式 AI (Generative AI)\n"
                "✓ 大數據 (Big Data)\n\n"
                "【已選場景主題 Selected Finance Scenario】\n"
                "✓ 風控 Risk Management\n"
                "✓ 普惠金融 Inclusive Finance",
            )

    # --- Slide 4 ---
    s4 = prs.slides[3]
    replace_in_slide(
        s4,
        "FINANO is an AI multi-agent engine",
        "FINANO is a Generative AI multi-agent engine",
    )
    for shape in s4.shapes:
        if hasattr(shape, "text") and "augments, not replaces" in shape.text:
            set_text(
                shape,
                "* FINANO augments, not replaces, a bank's existing KYC flow.\n"
                "* Aligns with BOCHK Challenge 2026: Innovation Leads the Future.",
            )

    # --- Slide 5: core engines ---
    s5 = prs.slides[4]
    mapping = {
        "8-question behavioral profile": (
            "FBTI: 8-question behavioral profile → 16 types, injected into every analysis."
        ),
        "5 parallel AI agents": (
            "MAFB: 5 parallel Generative AI agents score a fund; compliance agent vets output."
        ),
        "links every trade": (
            "Trade Replay: links every trade to its rationale, building a behavioral memory moat."
        ),
    }
    for shape in s5.shapes:
        if not hasattr(shape, "text"):
            continue
        for key, val in mapping.items():
            if key in shape.text:
                set_text(shape, val)

    # --- Slide 6 ---
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if hasattr(shape, "text"):
            if shape.text.strip() == "1.3 Value to BOCHK":
                set_text(shape, "1.3 Empowering BOCHK as an Innovation Pioneer")
            if shape.text.strip() == "WITH FINANO":
                set_text(shape, "WITH FINANO (GenAI)")
            if "A behavioral layer running" in shape.text:
                set_text(
                    shape,
                    "GenAI-powered behavioral layer running in parallel with KYC\n"
                    "Suitability-first matching, not return-first\n"
                    "Auto-logged rationale → satisfies \"record each recommendation's reason\"",
                )

    # --- Slide 10 ---
    s10 = prs.slides[9]
    for shape in s10.shapes:
        if hasattr(shape, "text") and "No incumbent combines" in shape.text:
            set_text(
                shape,
                "No incumbent combines behavioral profiling + Generative AI multi-agent "
                "analysis + trade memory — inside the bank's compliance perimeter.",
            )

    # --- Slide 14: prototype note ---
    s14 = prs.slides[13]
    for shape in s14.shapes:
        if hasattr(shape, "text") and "FINANO runs alongside KYC" in shape.text:
            set_text(
                shape,
                "FINANO runs alongside KYC — it closes the behavioral gap the questionnaire "
                "structurally cannot.\n\n"
                "Live prototype: React + FastAPI + LangGraph MAFB console (Docker-deployed demo).",
            )

    # --- Slide 18: pilot CTA ---
    s18 = prs.slides[17]
    for shape in s18.shapes:
        if hasattr(shape, "text") and "Success =" in shape.text:
            set_text(
                shape,
                "Success = measurable lift in suitability and compliance efficiency — "
                "proof we can scale to other GBA banks.\n\n"
                "Why we are here: Winning BOCHK Challenge 2026 is Step 1 to launch this "
                "PoC with BOCHK and unlock the competition's collaboration / internship pathway.",
            )
        if hasattr(shape, "text") and "What We" in shape.text:
            set_text(shape, "4.3 The Pilot: What We'll Prove")

    # --- Slide 21: cost structure ---
    s21 = prs.slides[20]
    for shape in s21.shapes:
        if hasattr(shape, "text") and "Compute" in shape.text and "MAFB" in shape.text:
            set_text(
                shape,
                "R&D — engineering team (largely fixed, already built)\n"
                "Compute — per MAFB analysis: 5 GenAI agents + compliance review "
                "(optimized via API routing & local LLM fallback)\n"
                "Data APIs — fund NAV / market data feeds\n"
                "Ops & hosting — servers, vector DB, maintenance",
            )

    # --- Slide 23: challenges content ---
    s23 = prs.slides[22]
    has_body = any(
        hasattr(s, "text") and "Challenge 1" in s.text for s in s23.shapes
    )
    if not has_body:
        left, top, width, height = 457200, 1600200, 8229600, 3200400
        box = s23.shapes.add_textbox(left, top, width, height)
        set_text(
            box,
            "Challenge 1 — Regulatory boundary: FINANO is decision-support, not licensed advice; "
            "compliance agent + disclaimers + human-in-the-loop.\n"
            "Challenge 2 — GenAI cost & latency: hybrid model routing, semantic caching, rule fallback.\n"
            "Challenge 3 — Data reliability: snapshot cache + throttled live feeds + graceful degradation.\n"
            "Challenge 4 — Bank integration: start with white-label API / RM console pilot, not full core replacement.",
        )

    prs.save(str(OUT))
    prs.save(str(COPY))
    print(f"Saved: {OUT}")
    print(f"Copied: {COPY}")


if __name__ == "__main__":
    main()
