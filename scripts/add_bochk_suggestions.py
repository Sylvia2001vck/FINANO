"""Add improvement suggestions below each slide — keep all original text untouched."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SRC = Path(
    r"c:\Users\30362\xwechat_files\wxid_iioe8mapqfa322_92bd\msg\file\2026-06\FINANO_Draft_0609.pptx"
)
OUT = SRC.with_name("FINANO_Draft_0609_with_suggestions.pptx")
COPY = Path(r"D:\FINANO\FINANO_Draft_0609_with_suggestions.pptx")

# Bilingual suggestions keyed by slide number (1-based)
SUGGESTIONS: dict[int, str] = {
    1: (
        "【改进建议】补全队员姓名与 HKUST 学号；在方案简介框写入约100字中英摘要，"
        "明确科技主题（生成式AI、大数据）与场景（风控、普惠金融）。"
        "演讲时先用1句话点出：FINANO 用 GenAI 多智能体解决银行「适当性+合規留痕」痛点。\n"
        "[Suggestion] Fill team roster & ~100-word bilingual summary; tag GenAI + Big Data "
        "and Risk Mgmt / Inclusive Finance. Open pitch with suitability + compliance angle."
    ),
    2: (
        "【改进建议】在官方模板勾选：生成式AI、大数据；场景至少选「风控」，可加「普惠金融」。"
        "右下角「量子計算」与 FINANO 不符，提交前请改为空白或删除该条（勿改上方官方选项文字）。"
        "本页是评委核对赛道的关键页，演讲30秒内说明所选主题即可。\n"
        "[Suggestion] Tick GenAI + Big Data; scenarios: Risk Mgmt (+ Inclusive Finance). "
        "Remove/replace the Quantum Computing line — it mismatches FINANO."
    ),
    3: (
        "【改进建议】章节页：口播衔接句示例——「Part 1 先回答 FINANO 是什么、为何适合中银香港」。"
        "保持原文标题不动，用一句话把听众带入痛点。\n"
        "[Suggestion] Section bridge: what FINANO is and why it fits BOCHK — do not rewrite slide title."
    ),
    4: (
        "【改进建议】在原文下方口播/增补时加入关键词 Generative AI（生成式AI），"
        "并呼应大赛主题 Innovation Leads the Future。"
        "强调「augment KYC」而非替代持牌投顾，避免合规误解。\n"
        "[Suggestion] When presenting, say 'Generative AI multi-agent engine'; stress decision-support, not licensed advice."
    ),
    5: (
        "【改进建议】口播时给三引擎命名：FBTI（画像）、MAFB（五路GenAI分析+合规）、Trade Replay（行为记忆）。"
        "建议在本节后插入1页系统截图/架构图（LangGraph 流程），满足大赛 Prototype 评分项。\n"
        "[Suggestion] Label engines FBTI / MAFB / Trade Replay; add a demo screenshot slide after this section."
    ),
    6: (
        "【改进建议】标题可口播为「赋能中银香港成为创新先驱」；WITH FINANO 一栏口播强调 "
        "GenAI-powered behavioral layer。对比表很好，演讲时先讲 TODAY 痛点再翻转。\n"
        "[Suggestion] Verbal title: Empowering BOCHK as Innovation Pioneer; highlight GenAI behavioral layer in WITH column."
    ),
    7: (
        "【改进建议】章节过渡：「Part 2 说明我们卖给谁——银行 B2B，中银香港是第一合作伙伴」。\n"
        "[Suggestion] Bridge to B2B buyer story; BOCHK as first partner."
    ),
    8: (
        "【改进建议】强化「本大赛 = PoC 第一步」与官网「中银香港合作/实习机会」的对应关系。"
        "Retail users 仅作验证层——评委问商业模式时用这个脚注回答。\n"
        "[Suggestion] Tie this competition to BOCHK PoC step 1 and official collaboration/internship pathway."
    ),
    9: (
        "【改进建议】数据很有说服力，演讲时只念 3 个数字：35.1T AUM、销售翻倍、20 家零售银行。"
        "保留脚注来源；若时间紧，SAM/SOM 用「1→3 家银行」一句话带过。\n"
        "[Suggestion] Present only 3 headline numbers; keep source footnote visible for credibility."
    ),
    10: (
        "【改进建议】增补口播词：加入 Generative AI multi-agent + trade memory。"
        "可补一张竞品对比表（智能投顾 / 基金超市 / 通用ChatGPT套壳 / FINANO 四列）。\n"
        "[Suggestion] Add 'Generative AI' in narration; optional competitor comparison table as extra slide."
    ),
    11: (
        "【改进建议】章节页：「Part 3 用监管与成本数据证明痛点真实存在」。\n"
        "[Suggestion] Set up regulatory + cost evidence in the next two slides."
    ),
    12: (
        "【改进建议】结构清晰。演讲强调：风险等级 ≠ 行为画像；机械匹配不足以满足 SFC 适当性要求。"
        "这是全篇道德与商业正当性基础，勿压缩时间。\n"
        "[Suggestion] Stress risk tier ≠ behavioral profile; mechanical matching fails suitability — keep full time here."
    ),
    13: (
        "【改进建议】HK$10.85M 罚款案例非常适合 1.5 分钟电梯演讲开场。"
        "三个指标（4.3M / 2x / 罚款）建议做口播重音，不必改原文。\n"
        "[Suggestion] Use HK$10.85M fine as elevator-pitch hook; emphasize three KPIs verbally."
    ),
    14: (
        "【改进建议】建议在本页后新增一页：Live Demo（MAFB 控制台、FBTI、合规 JSON 输出截图）。"
        "口播补一句：React + FastAPI + LangGraph，Docker 可部署。\n"
        "[Suggestion] Insert a demo screenshot slide after this one — critical for Prototype scoring."
    ),
    15: (
        "【改进建议】章节过渡：「Part 4 讲银行怎么买、我们怎么收钱」。\n"
        "[Suggestion] Transition to revenue model and BOCHK pilot."
    ),
    16: (
        "【改进建议】若视觉区「三种售卖方式」仍为空，请补：按席位 / 按次API / 年度白标许可。"
        "再次点题：本次大赛 = Pilot 第一步。\n"
        "[Suggestion] Fill three sales motions if blank on layout: per-seat, per-usage API, annual white-label."
    ),
    17: (
        "【改进建议】Y3 ARR 示意合理；提交前与团队统一 ACV 假设。"
        "口播强调：每增一家银行，边际服务成本几乎不变 → 高毛利 recurring。\n"
        "[Suggestion] Align ACV assumptions with team; narrate high-margin recurring revenue logic."
    ),
    18: (
        "【改进建议】在原文 Success 句之后口播补充："
        "「参加 BOCHK Challenge 2026 的目的，就是争取与中银香港落地 3 个月 A/B PoC」。"
        "明确成功指标：适当性提升%、合规记录耗时下降%、投诉率 proxy。\n"
        "[Suggestion] Verbal CTA: winning enables BOCHK PoC; define measurable pilot KPIs."
    ),
    19: (
        "【改进建议】章节过渡：「Part 5 证明银行端 ROI 与我们定价底线」。\n"
        "[Suggestion] Bridge to bank ROI and unit economics."
    ),
    20: (
        "【改进建议】若页面未列出四项银行收益，建议新增一页明细（不改本页原文）："
        "降投诉/罚款、降 RM 记录工时、提升交叉销售、提升客户留存。\n"
        "[Suggestion] If four gains are not visualized, add a detail slide — do not delete this slide's text."
    ),
    21: (
        "【改进建议】请工程组填写 HK$[___] 单次 MAFB 算力成本；"
        "口播补充：API routing + local LLM fallback 控制 GenAI token 成本。\n"
        "[Suggestion] Fill compute cost placeholder; mention hybrid routing & local fallback for GenAI cost control."
    ),
    22: (
        "【改进建议】提交前替换 HK$ xxxxxx / Xxx banks 为团队测算值；"
        "口播：固定成本为主，盈亏平衡看「签约银行数」。\n"
        "[Suggestion] Replace placeholder break-even numbers before submission."
    ),
    23: (
        "【改进建议】本页正文尚空，建议新增一页（保留本标题）写入四条挑战与对策："
        "①监管边界→合规节点+人工复核；②GenAI成本/延迟→缓存+规则降级；"
        "③数据可靠性→快照+节流；④银行集成→白标API/RM 控制台试点。\n"
        "[Suggestion] Add content slide: regulatory boundary, GenAI cost, data reliability, bank integration."
    ),
}


def add_suggestion_box(slide, text: str) -> None:
    left = Inches(0.35)
    top = Inches(6.05)
    width = Inches(12.6)
    height = Inches(1.35)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x8B, 0x00, 0x00)
    p.alignment = PP_ALIGN.LEFT
    # light background via shape fill
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE7)
    line = box.line
    line.color.rgb = RGBColor(0xCC, 0x66, 0x00)


def main() -> None:
    prs = Presentation(str(SRC))
    for idx, slide in enumerate(prs.slides, start=1):
        suggestion = SUGGESTIONS.get(idx)
        if suggestion:
            add_suggestion_box(slide, suggestion)
    prs.save(str(OUT))
    prs.save(str(COPY))
    print(f"Saved: {OUT}")
    print(f"Copied: {COPY}")


if __name__ == "__main__":
    main()
