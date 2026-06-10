#!/usr/bin/env python3
"""Export pitch-deck/index.html to an editable PowerPoint (.pptx)."""
from __future__ import annotations

import re
import sys
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
HTML_PATH = ROOT / "pitch-deck" / "index.html"
OUT_PATH = ROOT / "pitch-deck" / "FINANO-Pitch-Deck.pptx"
ASSETS = ROOT / "pitch-deck" / "assets"

BG = RGBColor(9, 9, 11)
WHITE = RGBColor(250, 250, 250)
MUTED = RGBColor(161, 161, 170)
ACCENT = RGBColor(212, 212, 216)

MARGIN_L = Inches(0.55)
MARGIN_T = Inches(0.45)
CONTENT_W = Inches(12.2)
MAX_BODY_H = Inches(6.2)


def clean(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_paragraph(text_frame, text: str, *, size: int = 14, bold: bool = False, color=WHITE, space_before: int = 6):
    if not text:
        return
    p = text_frame.add_paragraph() if text_frame.text else text_frame.paragraphs[0]
    if text_frame.text and text_frame.paragraphs[0].text:
        p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(space_before)
    p.level = 0


def add_text_block(slide, top, lines: list[tuple[str, int, bool, RGBColor]]):
    """lines: (text, font_size, bold, color)"""
    height = min(MAX_BODY_H, Inches(0.28 * max(len(lines), 1)))
    box = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for text, size, bold, color in lines:
        if not text:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_before = Pt(4 if size <= 14 else 8)
    return top + height + Inches(0.12)


def try_add_image(slide, rel_path: str, top, max_h=Inches(3.8)):
    rel = rel_path.lstrip("./")
    path = ASSETS / Path(rel).name if "/" not in rel else ROOT / "pitch-deck" / rel
    if not path.is_file():
        return top, False
    try:
        pic = slide.shapes.add_picture(str(path), MARGIN_L, top, width=CONTENT_W)
        ratio = pic.height / pic.width if pic.width else 1
        pic.width = int(CONTENT_W)
        pic.height = int(CONTENT_W * ratio)
        if pic.height > max_h:
            pic.height = int(max_h)
            pic.width = int(max_h / ratio)
        return top + Inches(pic.height / 914400) + Inches(0.1), True
    except Exception:
        return top, False


def extract_cards(section) -> list[str]:
    chunks: list[str] = []
    selectors = [
        ".toc-card",
        ".glass-card",
        ".focus-card",
        ".team-card",
        ".pricing-card",
        ".agenda-grid article",
    ]
    seen = set()
    for sel in selectors:
        for card in section.select(sel):
            bits: list[str] = []
            h = card.find(["h3", "h2"])
            if h:
                bits.append(clean(h.get_text()))
            role = card.select_one(".team-role")
            if role:
                bits.append(clean(role.get_text()))
            for p in card.find_all("p"):
                if p.find_parent(class_="team-copy") and p.get("class") and "team-role" in p.get("class", []):
                    continue
                t = clean(p.get_text())
                if t and t not in bits:
                    bits.append(t)
            for li in card.find_all("li"):
                t = clean(li.get_text())
                if t:
                    bits.append(f"• {t}")
            for span in card.select("span"):
                if span.find_parent("h3") or span.find_parent(".pricing-price"):
                    continue
                t = clean(span.get_text())
                if t and len(t) < 80 and t not in bits:
                    bits.append(t)
            block = "\n".join(bits)
            if block and block not in seen:
                seen.add(block)
                chunks.append(block)
    return chunks


def extract_lists(section) -> list[str]:
    items: list[str] = []
    for ul in section.select("ul, ol"):
        if ul.find_parent(class_=["glass-card", "toc-card", "team-card", "pricing-card"]):
            continue
        for li in ul.find_all("li", recursive=False):
            t = clean(li.get_text())
            if t:
                items.append(f"• {t}")
    return items


def extract_skill_matrix(section) -> list[str]:
    rows: list[str] = []
    for div in section.select(".skill-matrix div"):
        span = div.find("span")
        strong = div.find("strong")
        if span and strong:
            rows.append(f"• {clean(span.get_text())}: {clean(strong.get_text())}")
    return rows


def build_slide_content(section) -> tuple[str, list[str], list[str], list[tuple[str, str]]]:
    """Returns title, body lines, notes, images (rel, alt)."""
    title_parts: list[str] = []
    body: list[str] = []
    images: list[tuple[str, str]] = []

    for sel in [".eyebrow"]:
        el = section.select_one(sel)
        if el:
            body.append(f"[{clean(el.get_text())}]")

    for sel in [".slide-title", ".gradient-title", ".intro-title", ".metallic-title.mega"]:
        for el in section.select(sel):
            t = clean(el.get_text())
            if t:
                title_parts.append(t)

    if not title_parts:
        h2 = section.find("h2")
        if h2:
            title_parts.append(clean(h2.get_text()))

    title = title_parts[0] if title_parts else "Slide"
    if len(title_parts) > 1:
        body.extend(title_parts[1:])

    for el in section.select(".lead, .caption"):
        t = clean(el.get_text())
        if t:
            body.append(t)

    if section.select_one("video.intro-video"):
        body.append("[Embed intro video: pitch-deck/assets/finano_concept.mp4]")

    body.extend(extract_cards(section))
    body.extend(extract_lists(section))
    body.extend(extract_skill_matrix(section))

    for img in section.select("img.editable-image"):
        src = img.get("src", "")
        alt = img.get("alt", "")
        if "data:image" in src:
            continue
        images.append((src, alt))

    # Deduplicate body while preserving order
    seen = set()
    unique_body: list[str] = []
    for line in body:
        if line and line not in seen:
            seen.add(line)
            unique_body.append(line)

    return title, unique_body, [], images


def export_pptx(html_path: Path, out_path: Path) -> int:
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "html.parser")
    sections = soup.select("main#deck section.slide")
    if not sections:
        print("No slides found.", file=sys.stderr)
        return 1

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, section in enumerate(sections, start=1):
        slide = prs.slides.add_slide(blank)
        set_slide_bg(slide)

        title, body_lines, _, images = build_slide_content(section)
        top = MARGIN_T

        # Slide number
        num_box = slide.shapes.add_textbox(Inches(12.3), Inches(0.2), Inches(0.8), Inches(0.3))
        np = num_box.text_frame.paragraphs[0]
        np.text = f"{idx:02d}"
        np.font.size = Pt(10)
        np.font.color.rgb = MUTED
        np.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

        # Title
        tbox = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, Inches(0.9))
        tp = tbox.text_frame.paragraphs[0]
        tp.text = title[:200]
        tp.font.size = Pt(28 if len(title) < 60 else 22)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
        top += Inches(0.95)

        # Images (first one prominent; team avatars smaller)
        team_imgs = [i for i in images if "team-" in i[0]]
        other_imgs = [i for i in images if "team-" not in i[0]]

        for rel, alt in other_imgs[:2]:
            top, ok = try_add_image(slide, rel, top, max_h=Inches(3.5))
            if not ok and alt:
                body_lines.append(f"[Image: {alt}]")

        if team_imgs:
            body_lines.append("[Team photos included in HTML — see assets/team-*.png]")

        # Body text
        if body_lines:
            lines = []
            for block in body_lines[:24]:
                if "\n" in block:
                    for sub in block.split("\n"):
                        if sub.strip():
                            lines.append((sub.strip(), 13, sub.startswith("•"), MUTED if sub.startswith("•") else WHITE))
                else:
                    bold = block.startswith("[") and block.endswith("]")
                    lines.append((block, 12 if bold else 14, bold, ACCENT if bold else WHITE))
            add_text_block(slide, top, lines)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {len(sections)} slides -> {out_path}")
    return 0


if __name__ == "__main__":
    html = Path(sys.argv[1]) if len(sys.argv) > 1 else HTML_PATH
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else OUT_PATH
    raise SystemExit(export_pptx(html, out))
